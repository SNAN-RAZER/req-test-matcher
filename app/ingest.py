"""Ingest PDF, Excel, PowerPoint (and nested attachments) without leaving the LAN."""

from __future__ import annotations

import shutil
import zipfile
from dataclasses import dataclass, field
from pathlib import Path

from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from app.config import settings

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".bmp", ".gif"}
SUPPORTED = {".pdf", ".xlsx", ".xlsm", ".xls", ".pptx", ".ppt", ".docx", ".doc", ".txt", ".md", ".csv"} | IMAGE_EXTS
OOXML = {".xlsx", ".xlsm", ".pptx", ".docx"}


@dataclass
class ParsedDoc:
    path: Path
    source_name: str
    parent_name: str
    text: str
    kind: str
    images: list[Path] = field(default_factory=list)
    modality: str = "semantic"
    nested: list["ParsedDoc"] = field(default_factory=list)


def extract_embedded_from_ooxml(path: Path, out_dir: Path) -> list[Path]:
    found: list[Path] = []
    if path.suffix.lower() not in OOXML:
        return found
    try:
        with zipfile.ZipFile(path) as zf:
            for name in zf.namelist():
                lower = name.lower()
                if not any(seg in lower for seg in ("/embeddings/", "/oleobjects/", "embeddings/")):
                    continue
                if name.endswith("/"):
                    continue
                data = zf.read(name)
                if not data:
                    continue
                dest = out_dir / Path(name).name
                dest.parent.mkdir(parents=True, exist_ok=True)
                if dest.exists():
                    dest = out_dir / f"{path.stem}_{Path(name).name}"
                dest.write_bytes(data)
                found.append(dest)
    except zipfile.BadZipFile:
        return found
    return found


def extract_embedded_from_pdf(path: Path, out_dir: Path) -> list[Path]:
    found: list[Path] = []
    if path.suffix.lower() != ".pdf":
        return found
    try:
        reader = PdfReader(str(path))
        attachments = getattr(reader, "attachments", None) or {}
        for filename, payload in attachments.items():
            data = payload[0] if isinstance(payload, list) else payload
            if isinstance(data, str):
                data = data.encode("utf-8", errors="replace")
            dest = out_dir / Path(str(filename)).name
            dest.write_bytes(data)
            found.append(dest)
    except Exception:
        return found
    return found


def _text_from_pdf(path: Path) -> str:
    try:
        reader = PdfReader(str(path))
        pages = []
        for page in reader.pages:
            pages.append(page.extract_text() or "")
        return "\n\n".join(pages)
    except Exception as exc:
        return f"[pdf parse error: {exc}]"


def _text_from_xlsx(path: Path) -> str:
    try:
        wb = load_workbook(path, data_only=True, read_only=True)
        chunks: list[str] = []
        for sheet in wb.worksheets:
            chunks.append(f"# Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = ["" if c is None else str(c).strip() for c in row]
                if any(cells):
                    chunks.append(" | ".join(cells))
        wb.close()
        return "\n".join(chunks)
    except Exception as exc:
        return f"[excel parse error: {exc}]"


def _text_from_pptx(path: Path) -> str:
    try:
        prs = Presentation(str(path))
        chunks: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            chunks.append(f"# Slide {i}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        chunks.append(" | ".join(cells))
        return "\n".join(chunks)
    except Exception as exc:
        return f"[pptx parse error: {exc}]"


def _text_plain(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="replace")
    except Exception:
        return path.read_bytes().decode("utf-8", errors="replace")


def _try_docling(path: Path) -> str | None:
    try:
        from docling.document_converter import DocumentConverter

        converter = DocumentConverter()
        result = converter.convert(str(path))
        md = result.document.export_to_markdown()
        return md if md and md.strip() else None
    except Exception:
        return None


def extract_text(path: Path) -> tuple[str, str]:
    suffix = path.suffix.lower()
    docling_text = _try_docling(path)
    if docling_text:
        return docling_text, "docling"
    if suffix == ".pdf":
        return _text_from_pdf(path), "pypdf"
    if suffix == ".csv":
        return _text_plain(path).replace(",", " | "), "csv"
    if suffix in {".xlsx", ".xlsm", ".xls"}:
        return _text_from_xlsx(path), "openpyxl"
    if suffix in {".pptx", ".ppt"}:
        return _text_from_pptx(path), "python-pptx"
    if suffix in {".txt", ".md"}:
        return _text_plain(path), "plain"
    if suffix == ".docx":
        try:
            with zipfile.ZipFile(path) as zf:
                xml = zf.read("word/document.xml").decode("utf-8", errors="replace")
            import re

            texts = re.findall(r"<w:t[^>]*>(.*?)</w:t>", xml)
            return "\n".join(texts), "docx-xml"
        except Exception as exc:
            return f"[docx parse error: {exc}]", "error"
    return _text_plain(path), "fallback"


def extract_images(path: Path, out_dir: Path) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    found: list[Path] = []
    suffix = path.suffix.lower()
    if suffix in IMAGE_EXTS:
        dest = out_dir / path.name
        if path.resolve() != dest.resolve():
            shutil.copy2(path, dest)
            found.append(dest)
        else:
            found.append(path)
        return found
    if suffix == ".pdf":
        try:
            reader = PdfReader(str(path))
            for i, page in enumerate(reader.pages):
                for j, img in enumerate(getattr(page, "images", []) or []):
                    name = getattr(img, "name", f"p{i}_{j}.bin")
                    dest = out_dir / f"{path.stem}_p{i}_{j}_{Path(str(name)).name}"
                    data = getattr(img, "data", None)
                    if data:
                        dest.write_bytes(data)
                        found.append(dest)
        except Exception:
            return found
        return found
    if suffix in {".pptx", ".docx", ".xlsx"}:
        try:
            with zipfile.ZipFile(path) as zf:
                for name in zf.namelist():
                    lower = name.lower()
                    if "/media/" not in lower and "/embeddings/" not in lower:
                        continue
                    if Path(name).suffix.lower() not in IMAGE_EXTS:
                        continue
                    dest = out_dir / f"{path.stem}_{Path(name).name}"
                    dest.write_bytes(zf.read(name))
                    found.append(dest)
        except zipfile.BadZipFile:
            return found
    return found


def classify_modality(path: Path, text: str, images: list[Path]) -> str:
    suffix = path.suffix.lower()
    if suffix in IMAGE_EXTS:
        return "vision"
    if suffix in {".xlsx", ".xlsm", ".xls", ".csv"}:
        return "table"
    tableish = text.count("|") >= 8 or text.lower().count("sheet:") >= 1
    if images and tableish:
        return "mixed"
    if images and not (text or "").strip():
        return "vision"
    if images:
        return "mixed"
    if tableish:
        return "table"
    return "semantic"


def ingest_file(path: Path, parent: str = "", depth: int = 0, max_depth: int = 4) -> ParsedDoc:
    embed_dir = settings.work_dir / "embedded" / path.stem
    embed_dir.mkdir(parents=True, exist_ok=True)
    image_dir = settings.work_dir / "images" / path.stem
    nested_paths = extract_embedded_from_ooxml(path, embed_dir) + extract_embedded_from_pdf(path, embed_dir)
    images = extract_images(path, image_dir)
    text, kind = extract_text(path)
    nested: list[ParsedDoc] = []
    if depth < max_depth:
        for child in nested_paths:
            if child.suffix.lower() in SUPPORTED or child.suffix.lower() in OOXML:
                nested.append(ingest_file(child, parent=path.name, depth=depth + 1, max_depth=max_depth))
            else:
                sniff = child.read_bytes()[:8]
                if sniff.startswith(b"PK"):
                    try:
                        nested.append(ingest_file(child, parent=path.name, depth=depth + 1, max_depth=max_depth))
                    except Exception:
                        nested.append(
                            ParsedDoc(
                                path=child,
                                source_name=child.name,
                                parent_name=path.name,
                                text=f"[unsupported embedded blob: {child.name}]",
                                kind="blob",
                                modality="nested",
                            )
                        )
                else:
                    nested.append(
                        ParsedDoc(
                            path=child,
                            source_name=child.name,
                            parent_name=path.name,
                            text=f"[embedded file stored: {child.name}]",
                            kind="blob",
                            modality="nested",
                        )
                    )
    combined = text
    for n in nested:
        combined += f"\n\n--- Embedded: {n.source_name} (inside {path.name}) ---\n{n.text}"
        images.extend(n.images)
    modality = classify_modality(path, combined, images)
    return ParsedDoc(
        path=path,
        source_name=path.name,
        parent_name=parent,
        text=combined,
        kind=kind,
        images=images,
        modality=modality,
        nested=nested,
    )


def flatten_docs(docs: list[ParsedDoc]) -> list[ParsedDoc]:
    out: list[ParsedDoc] = []
    for d in docs:
        out.append(d)
        out.extend(flatten_docs(d.nested))
    return out
